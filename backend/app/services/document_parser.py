import logging
import os

import fitz
from docx import Document as DocxDocument
from openpyxl import load_workbook

logger = logging.getLogger(__name__)


def parse_text(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8") as f:
        return f.read()


def parse_pdf(filepath: str) -> str:
    text_parts = []
    with fitz.open(filepath) as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n".join(text_parts)


def parse_docx(filepath: str) -> str:
    doc = DocxDocument(filepath)
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)


def parse_xlsx(filepath: str) -> str:
    wb = load_workbook(filepath, read_only=True, data_only=True)
    result_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        result_parts.append(f"## {sheet_name}")
        rows = list(ws.iter_rows(values_only=True))
        if rows:
            headers = [str(c) if c is not None else "" for c in rows[0]]
            result_parts.append("| " + " | ".join(headers) + " |")
            result_parts.append("| " + " | ".join(["---"] * len(headers)) + " |")
            for row in rows[1:]:
                cells = [str(c) if c is not None else "" for c in row]
                result_parts.append("| " + " | ".join(cells) + " |")
        result_parts.append("")
    wb.close()
    return "\n".join(result_parts)


def parse_xls(filepath: str) -> str:
    import xlrd

    wb = xlrd.open_workbook(filepath)
    result_parts = []
    for idx in range(wb.nsheets):
        ws = wb.sheet_by_index(idx)
        sheet_name = ws.name
        result_parts.append(f"## {sheet_name}")
        if ws.nrows > 0:
            headers = [str(ws.cell_value(0, c)) for c in range(ws.ncols)]
            result_parts.append("| " + " | ".join(headers) + " |")
            result_parts.append("| " + " | ".join(["---"] * len(headers)) + " |")
            for r in range(1, ws.nrows):
                cells = [str(ws.cell_value(r, c)) for c in range(ws.ncols)]
                result_parts.append("| " + " | ".join(cells) + " |")
        result_parts.append("")
    return "\n".join(result_parts)


def parse_pptx(filepath: str) -> str:
    from pptx import Presentation

    prs = Presentation(filepath)
    result_parts = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_title = ""
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if shape.is_placeholder and shape.placeholder_format.type == 1:
                        slide_title = text
                    else:
                        texts.append(text)

            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    result_parts.append("\n".join(rows))
                    result_parts.append("")

        result_parts.append(f"# {slide_title or f'幻灯片 {idx}'}")
        if texts:
            result_parts.append("\n".join(texts))
        result_parts.append("")

    return "\n\n".join(result_parts)


def parse_markdown(filepath: str) -> str:
    return parse_text(filepath)


def convert_to_markdown(filepath: str) -> str:
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".txt":
        return parse_text(filepath)
    elif ext == ".md":
        return parse_markdown(filepath)
    elif ext == ".pdf":
        content = parse_pdf(filepath)
        return _format_as_markdown(content)
    elif ext == ".docx":
        content = parse_docx(filepath)
        return _format_as_markdown(content)
    elif ext == ".xlsx":
        return parse_xlsx(filepath)
    elif ext == ".xls":
        return parse_xls(filepath)
    elif ext == ".pptx":
        return parse_pptx(filepath)
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def extract_text(filepath: str) -> str:
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return parse_pdf(filepath)
    elif ext == ".docx":
        return parse_docx(filepath)
    elif ext == ".xlsx":
        return parse_xlsx(filepath)
    elif ext == ".xls":
        return parse_xls(filepath)
    else:
        return parse_text(filepath)


def _format_as_markdown(text: str) -> str:
    lines = text.strip().split("\n")
    result = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            result.append("")
        elif len(stripped) < 80 and stripped.endswith(("。", "：", "）", ")")):
            result.append(f"### {stripped}")
        else:
            result.append(stripped)
    return "\n\n".join(result)
